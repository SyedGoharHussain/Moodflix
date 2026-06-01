"""
MoodFlix — presentation builder.
Generates MoodFlix_Presentation.pptx : dark cinematic theme, red accent,
real app screenshots, simple plain-English copy.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

MEDIA = "/tmp/moodflix_docx/word/media"
OUT = "/Users/usman/Desktop/moodflix/MoodFlix_Presentation.pptx"

# ---- palette (no gold, no blue) ----
BG    = RGBColor(0x0E, 0x0E, 0x10)
CARD  = RGBColor(0x19, 0x19, 0x1D)
CARD2 = RGBColor(0x21, 0x21, 0x26)
RED   = RGBColor(0xE5, 0x09, 0x14)
REDLT = RGBColor(0xFF, 0x5A, 0x63)
WHITE = RGBColor(0xF5, 0xF5, 0xF6)
GRAY  = RGBColor(0xA6, 0xA6, 0xAC)
DIM   = RGBColor(0x77, 0x77, 0x7E)
LINE  = RGBColor(0x2C, 0x2C, 0x33)

FONT  = "Segoe UI"
FONTL = "Segoe UI Light"
FONTB = "Segoe UI Semibold"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    r._element.addprevious(r._element)
    return s


def _set(p, text, size, color, bold=False, font=FONT, italic=False, spacing=None):
    p.text = text if text != "" else " "
    # A '\n' splits into multiple runs (with <a:br/>); format EVERY run.
    for r in p.runs:
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
        if spacing is not None:
            r.font._rPr.set("spc", str(int(spacing * 100)))
    return p


def textbox(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        _set(p, ln["t"], ln["s"], ln.get("c", WHITE), ln.get("b", False),
             ln.get("f", FONT), ln.get("i", False), ln.get("sp"))
        if ln.get("sa") is not None: p.space_after = Pt(ln["sa"])
        if ln.get("sb") is not None: p.space_before = Pt(ln["sb"])
        if ln.get("lh") is not None: p.line_spacing = ln["lh"]
    return tb


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if shadow:
        _shadow(sp)
    return sp


def _shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '90000', 'dist': '50000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '55000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)


def kicker(s, text):
    rect(s, 0.7, 0.62, 0.34, 0.07, fill=RED)
    textbox(s, 1.12, 0.5, 9, 0.4, [{"t": text.upper(), "s": 13, "c": REDLT, "b": True, "sp": 2.5}])


def title(s, text, size=33):
    textbox(s, 0.68, 0.85, 11.9, 1.0, [{"t": text, "s": size, "c": WHITE, "b": True}])


def footer(s, n):
    textbox(s, 0.7, 7.04, 4, 0.3, [{"t": "MoodFlix  ·  AI Mood Detector for Movies", "s": 9, "c": DIM}])
    textbox(s, 11.4, 7.04, 1.23, 0.3, [{"t": f"{n:02d}", "s": 9, "c": DIM, "align": PP_ALIGN.RIGHT}])


def fit_image(s, path, bl, bt, bw, bh, border=True, caption=None):
    iw, ih = Image.open(path).size
    ir, br = iw / ih, bw / bh
    if ir > br:
        w = bw; h = bw / ir
    else:
        h = bh; w = bh * ir
    l = bl + (bw - w) / 2; t = bt + (bh - h) / 2
    # backing frame for depth
    pad = 0.05
    fr = rect(s, l - pad, t - pad, w + 2 * pad, h + 2 * pad, fill=CARD2, line=LINE, lw=1, shadow=True)
    pic = s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    if caption:
        textbox(s, bl, bt + bh + 0.04, bw, 0.32,
                [{"t": caption, "s": 10.5, "c": GRAY, "align": PP_ALIGN.CENTER}])
    return pic


def img(name):
    return os.path.join(MEDIA, name)


# ============================================================ 1. COVER
s = slide()
# right-side hero screenshot
fit_image(s, img("image10.png"), 8.55, 0.0, 4.78, 7.5, border=False)
# left dark gradient panel over nothing (text sits on bg)
rect(s, 0, 0, 8.7, 7.5, fill=BG)
rect(s, 0.9, 1.05, 0.5, 0.09, fill=RED)
textbox(s, 0.86, 1.35, 7.6, 2.4, [
    {"t": "MoodFlix", "s": 78, "c": WHITE, "b": True, "sa": 2},
    {"t": "Netflix, but it reads the room.", "s": 22, "c": REDLT, "b": True, "f": FONT, "sa": 6},
    {"t": "An AI movie & TV recommendation engine that\npicks films from how you feel.", "s": 15, "c": GRAY, "lh": 1.25},
])
# info block
rect(s, 0.9, 4.55, 7.2, 0.018, fill=LINE)
textbox(s, 0.9, 4.75, 7.4, 2.4, [
    {"t": "AI SEMESTER PROJECT  ·  FINAL PRESENTATION", "s": 11.5, "c": REDLT, "b": True, "sp": 1.5, "sa": 4},
    {"t": "BS IT — VI B", "s": 13, "c": GRAY, "sa": 12},
    {"t": "Muhammad Usman    ·    230966", "s": 14.5, "c": WHITE, "b": True, "sa": 3},
    {"t": "Gohar Husain         ·    230930", "s": 14.5, "c": WHITE, "b": True, "sa": 3},
    {"t": "Ahmad Amir            ·    230948", "s": 14.5, "c": WHITE, "b": True, "sa": 10},
    {"t": "Submitted to:  Ms. Warda Aslam", "s": 13, "c": GRAY},
])

# ============================================================ 2. PROBLEM
s = slide(); kicker(s, "The Problem & Our Idea"); title(s, "A genre is not a feeling")
textbox(s, 0.7, 1.95, 6.0, 4.6, [
    {"t": "Most movie apps make you pick a genre. But a genre does not\nmatch your mood.", "s": 16, "c": GRAY, "lh": 1.3, "sa": 16},
    {"t": "•  On a Sunday you might want something cozy.", "s": 15, "c": WHITE, "sa": 8},
    {"t": "•  After a long week you might want something light.", "s": 15, "c": WHITE, "sa": 8},
    {"t": "•  Late at night you might want something dark or scary.", "s": 15, "c": WHITE, "sa": 18},
    {"t": "Our idea", "s": 13, "c": REDLT, "b": True, "sp": 1.5, "sa": 5},
    {"t": "You just type how you feel. The AI reads it, finds your mood,\nand shows movies that fit — and tells you why.", "s": 16.5, "c": WHITE, "b": True, "lh": 1.3},
])
fit_image(s, img("image16.png"), 7.55, 1.5, 5.1, 5.2,
          caption="The app detecting a “relaxed” mood from your text")
footer(s, 2)

# ============================================================ 3. HOW IT WORKS
s = slide(); kicker(s, "How It Works"); title(s, "From a sentence to the right movie")
flow = [
    ("1", "You type", "“had a long week,\nwant something cozy”"),
    ("2", "AI model reads it", "DistilBERT + TF-IDF\nunderstand the words"),
    ("3", "Mood detected", "e.g.  relaxed\n(with a confidence %)"),
    ("4", "Movies returned", "matching titles\nfrom TMDB"),
]
bx, by, bw, bh, gap = 0.7, 2.35, 2.78, 2.0, 0.45
for i, (n, h, d) in enumerate(flow):
    x = bx + i * (bw + gap)
    c = rect(s, x, by, bw, bh, fill=CARD, line=(RED if i in (1,) else LINE),
             lw=(1.6 if i == 1 else 1), shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    textbox(s, x, by + 0.22, bw, bh - 0.3, [
        {"t": n, "s": 15, "c": RED, "b": True, "align": PP_ALIGN.CENTER, "sa": 6},
        {"t": h, "s": 16, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER, "sa": 6},
        {"t": d, "s": 12, "c": GRAY, "align": PP_ALIGN.CENTER, "lh": 1.15},
    ])
    if i < 3:
        textbox(s, x + bw - 0.02, by + bh/2 - 0.35, gap + 0.04, 0.7,
                [{"t": "→", "s": 26, "c": RED, "b": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.7, 4.95, 12, 0.4, [{"t": "The 16 moods our AI can detect", "s": 13, "c": REDLT, "b": True, "sp": 1.2}])
moods = "happy  ·  sad  ·  lonely  ·  romantic  ·  excited  ·  relaxed  ·  stressed  ·  dark  ·  emotional  ·  mind-bending  ·  curious  ·  nostalgic  ·  motivated  ·  adventurous  ·  wholesome  ·  scared"
textbox(s, 0.7, 5.35, 11.95, 1.0, [{"t": moods, "s": 15.5, "c": WHITE, "lh": 1.45}])
footer(s, 3)

# ============================================================ 4. THE AI MODELS
s = slide(); kicker(s, "The AI Brain"); title(s, "Two AI models, working as a team")
cards = [
    ("MODEL A", "TF-IDF + Logistic Regression", REDLT, [
        "A classic, fast text model.", "Counts important words and",
        "picks the closest mood.", "", "Super fast (~5 ms), but weaker", "on slang and emojis."]),
    ("MODEL B  ·  MAIN", "DistilBERT", RED, [
        "A small, smart language model", "(a lighter version of BERT).",
        "We re-trained it to understand", "moods from real sentences.", "",
        "Handles slang, typos & emojis."]),
    ("LIVE SYSTEM", "The Ensemble", REDLT, [
        "We blend both models:", "65% DistilBERT + 35% TF-IDF.",
        "", "If both agree, confidence goes", "up. Emojis & keywords act as", "a safety net."]),
]
cw, ch, cgap = 3.86, 4.35, 0.3
cx0 = 0.7
for i, (k, h, kc, body) in enumerate(cards):
    x = cx0 + i * (cw + cgap)
    main = (i == 1)
    rect(s, x, 2.1, cw, ch, fill=(CARD2 if main else CARD),
         line=(RED if main else LINE), lw=(1.8 if main else 1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    textbox(s, x + 0.32, 2.4, cw - 0.6, 0.4, [{"t": k, "s": 11.5, "c": kc, "b": True, "sp": 1.5}])
    textbox(s, x + 0.32, 2.78, cw - 0.6, 0.7, [{"t": h, "s": 21, "c": WHITE, "b": True}])
    rect(s, x + 0.32, 3.5, 0.7, 0.045, fill=RED)
    textbox(s, x + 0.32, 3.72, cw - 0.6, ch - 1.7,
            [{"t": ln, "s": 13.5, "c": (WHITE if ln else GRAY), "lh": 1.2, "sa": 1} for ln in body])
footer(s, 4)

# ============================================================ 5. TRAINING
s = slide(); kicker(s, "How We Trained It"); title(s, "Teaching the model what moods look like")
textbox(s, 0.7, 1.95, 6.05, 5, [
    {"t": "The data", "s": 13, "c": REDLT, "b": True, "sp": 1.2, "sa": 6},
    {"t": "We built one big, balanced dataset of about 96,000\nlabelled sentences across all 16 moods.", "s": 15.5, "c": WHITE, "lh": 1.3, "sa": 7},
    {"t": "It mixes our own hand-written examples with two public\nemotion datasets, plus generated samples for rare moods.", "s": 14, "c": GRAY, "lh": 1.3, "sa": 16},
    {"t": "The training", "s": 13, "c": REDLT, "b": True, "sp": 1.2, "sa": 6},
    {"t": "•  Split: 90% to learn (86,400) · 10% to test (9,600)", "s": 14.5, "c": WHITE, "sa": 7},
    {"t": "•  3 rounds (epochs), batch size 32, fp16 fast mode", "s": 14.5, "c": WHITE, "sa": 7},
    {"t": "•  Loss dropped from 2.83 → 0.12 — it learned well", "s": 14.5, "c": WHITE, "sa": 7},
    {"t": "•  Finished in ~10 minutes on an RTX 4060 Ti GPU", "s": 14.5, "c": WHITE},
])
# clean preprocessing note card
rect(s, 7.15, 2.0, 5.45, 4.5, fill=CARD, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
textbox(s, 7.5, 2.35, 4.8, 4, [
    {"t": "BEFORE TRAINING WE CLEAN THE TEXT", "s": 12, "c": REDLT, "b": True, "sp": 1.2, "sa": 12},
    {"t": "✓  make lowercase", "s": 15, "c": WHITE, "sa": 9},
    {"t": "✓  remove links, @tags, #hashtags", "s": 15, "c": WHITE, "sa": 9},
    {"t": "✓  remove punctuation", "s": 15, "c": WHITE, "sa": 9},
    {"t": "✓  keep numbers (“10/10”, “5 stars”)", "s": 15, "c": WHITE, "sa": 16},
    {"t": "DistilBERT reads the raw text, so it keeps\nemojis, CAPS and slang as extra clues.", "s": 13, "c": GRAY, "i": True, "lh": 1.3},
])
footer(s, 5)

# ============================================================ 6. RESULTS
s = slide(); kicker(s, "Results"); title(s, "How well does it work?")
stats = [
    ("TF-IDF + LogReg", "87.21%", "accuracy  ·  very fast", False),
    ("DistilBERT (main)", "90.19%", "accuracy  ·  macro-F1 0.90", True),
    ("Live Ensemble", "91.4%", "real-world hits (32 / 35)", False),
]
cw, ch, gap = 3.86, 1.95, 0.3
for i, (h, big, sub, main) in enumerate(stats):
    x = 0.7 + i * (cw + gap)
    rect(s, x, 1.95, cw, ch, fill=(CARD2 if main else CARD),
         line=(RED if main else LINE), lw=(1.8 if main else 1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    textbox(s, x, 2.18, cw, 0.4, [{"t": h.upper(), "s": 12, "c": (REDLT if main else GRAY), "b": True, "align": PP_ALIGN.CENTER, "sp": 1}])
    textbox(s, x, 2.5, cw, 0.95, [{"t": big, "s": 46, "c": (RED if main else WHITE), "b": True, "align": PP_ALIGN.CENTER}])
    textbox(s, x, 3.45, cw, 0.4, [{"t": sub, "s": 12.5, "c": GRAY, "align": PP_ALIGN.CENTER}])
textbox(s, 0.7, 4.12, 11.9, 0.4,
        [{"t": "Live test on 35 mixed real inputs — formal, slang, Gen Z, emojis and typos:", "s": 13.5, "c": GRAY}])
fit_image(s, img("image9.png"), 0.7, 4.5, 11.95, 2.35,
          caption="Smoke test — the ensemble correctly classified 32 of 35 inputs (91.4%)")
footer(s, 6)

# ============================================================ 7. EXPLAINABLE AI
s = slide(); kicker(s, "Explainable AI"); title(s, "It shows WHY it chose your mood")
fit_image(s, img("image19.png"), 0.7, 1.62, 4.0, 4.9,
          caption="Real /api/analyze-mood response with word weights")
textbox(s, 5.1, 2.05, 7.4, 4.8, [
    {"t": "No black box", "s": 13, "c": REDLT, "b": True, "sp": 1.2, "sa": 8},
    {"t": "For every prediction we use LIME — a tool that shows\nwhich words pushed the AI toward a mood.", "s": 16, "c": WHITE, "lh": 1.3, "sa": 16},
    {"t": "In the example on the left:", "s": 14, "c": GRAY, "sa": 8},
    {"t": "•  detected mood:  relaxed  (96% sure)", "s": 15, "c": WHITE, "sa": 8},
    {"t": "•  the word “cozy” had the biggest positive weight", "s": 15, "c": WHITE, "sa": 8},
    {"t": "•  “tough” pulled slightly the other way", "s": 15, "c": WHITE, "sa": 18},
    {"t": "The app highlights these words for the user, so the\nresult feels trustworthy instead of random.", "s": 15, "c": GRAY, "lh": 1.3},
])
footer(s, 7)

# ============================================================ 8. APP 1
s = slide(); kicker(s, "Inside the App  ·  1"); title(s, "A clean, cinematic interface")
fit_image(s, img("image10.png"), 0.7, 1.7, 5.85, 5.05, caption="Landing page — “Netflix, but it reads the room.”")
fit_image(s, img("image11.png"), 6.85, 1.7, 5.78, 5.05, caption="Sign in — email / password or Google (Firebase)")
footer(s, 8)

# ============================================================ 9. APP 2
s = slide(); kicker(s, "Inside the App  ·  2"); title(s, "Mood detection in action")
fit_image(s, img("image16.png"), 0.7, 1.65, 4.35, 5.15, caption="“Because you feel relaxed”")
fit_image(s, img("image17.png"), 5.05, 1.65, 4.35, 5.15, caption="“Because you feel excited”")
textbox(s, 9.7, 2.1, 3.1, 4.6, [
    {"t": "Same app,\ndifferent feeling.", "s": 20, "c": WHITE, "b": True, "lh": 1.2, "sa": 14},
    {"t": "You type a sentence, the AI shows the detected mood with a\nconfidence bar, then fills the\npage with matching movies.", "s": 14.5, "c": GRAY, "lh": 1.3, "sa": 14},
    {"t": "Change how you feel and the\nwhole page changes with you.", "s": 14.5, "c": REDLT, "b": True, "lh": 1.3},
])
footer(s, 9)

# ============================================================ 10. APP 3
s = slide(); kicker(s, "Inside the App  ·  3"); title(s, "Every pick is explained")
fit_image(s, img("image15.png"), 0.7, 1.7, 5.85, 5.05, caption="Movie page — “Why this was recommended”")
fit_image(s, img("image14.png"), 6.85, 1.7, 5.78, 5.05, caption="Search — find any title fast")
footer(s, 10)

# ============================================================ 11. TECH STACK
s = slide(); kicker(s, "Tech Stack"); title(s, "What we built it with")
groups = [
    ("FRONTEND", ["React 18 + Vite", "Tailwind CSS (dark theme)", "Framer Motion animations", "Zustand state", "React Router"]),
    ("BACKEND", ["Python 3 + Flask", "REST API (/api/analyze-mood)", "Rate limiting middleware", "JSON responses"]),
    ("AI / ML", ["scikit-learn (TF-IDF)", "PyTorch + DistilBERT", "LIME explainability", "pandas, joblib"]),
    ("SERVICES", ["TMDB API (movies & TV)", "Firebase Auth + Firestore", "Google sign-in", "JWT sessions"]),
]
cw, ch, gap = 2.92, 4.4, 0.27
for i, (h, items) in enumerate(groups):
    x = 0.7 + i * (cw + gap)
    rect(s, x, 2.1, cw, ch, fill=CARD, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, 2.1, cw, 0.09, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x + 0.28, 2.4, cw - 0.5, 0.4, [{"t": h, "s": 13.5, "c": REDLT, "b": True, "sp": 1.5}])
    textbox(s, x + 0.28, 2.95, cw - 0.5, ch - 1,
            [{"t": "•  " + it, "s": 14, "c": WHITE, "lh": 1.2, "sa": 9} for it in items])
footer(s, 11)

# ============================================================ 12. REPOSITORY
s = slide(); kicker(s, "Our Repository"); title(s, "How the project is organised")
fit_image(s, img("image1.png"), 7.4, 0.55, 5.4, 6.85, border=True)
textbox(s, 0.7, 2.0, 6.3, 5, [
    {"t": "One clean repo, two halves.", "s": 18, "c": WHITE, "b": True, "sa": 16},
    {"t": "client/", "s": 16, "c": RED, "b": True, "sa": 3},
    {"t": "The React frontend — pages, components, services\nand the dark UI.", "s": 14, "c": GRAY, "lh": 1.25, "sa": 14},
    {"t": "server/", "s": 16, "c": RED, "b": True, "sa": 3},
    {"t": "The Flask backend and all the AI code in server/ml/ —\ntraining scripts, the dataset, and the saved models\n(mood_classifier.pkl + the DistilBERT checkpoint).", "s": 14, "c": GRAY, "lh": 1.25, "sa": 14},
    {"t": "Helper docs (README, RUNME) explain how to run it\nin a few commands.", "s": 13.5, "c": DIM, "lh": 1.25},
])
footer(s, 12)

# ============================================================ 13. HARDWARE STORY
s = slide(); kicker(s, "The Hardware Story"); title(s, "The Mac couldn’t finish it — the PC did")
two = [
    ("MacBook Air  M2", DIM, [
        ("Chip", "Apple M2 (MPS)"),
        ("DistilBERT job", "ran for hours"),
        ("Result", "overheated &\nkilled mid-run"),
        ("Saved model", "nothing"),
    ], "App fell back to the weaker\nTF-IDF model only."),
    ("Desktop  ·  RTX 4060 Ti", RED, [
        ("GPU", "RTX 4060 Ti (CUDA)"),
        ("DistilBERT job", "fp16 fast mode"),
        ("Result", "finished in\n~10 minutes"),
        ("Saved model", "90.19% accuracy"),
    ], "Trained the real model. It now\nruns on the Mac for inference."),
]
cw, gap = 5.85, 0.5
for i, (h, hc, rows, note) in enumerate(two):
    x = 0.7 + i * (cw + gap)
    main = i == 1
    rect(s, x, 2.0, cw, 4.45, fill=(CARD2 if main else CARD),
         line=(RED if main else LINE), lw=(1.8 if main else 1),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    textbox(s, x + 0.4, 2.3, cw - 0.8, 0.5, [{"t": h, "s": 19, "c": (WHITE if main else GRAY), "b": True}])
    rect(s, x + 0.4, 2.85, 0.8, 0.045, fill=hc)
    yy = 3.1
    for k, v in rows:
        textbox(s, x + 0.4, yy, 2.1, 0.5, [{"t": k, "s": 13, "c": GRAY}])
        textbox(s, x + 2.5, yy, cw - 2.9, 0.6, [{"t": v, "s": 14, "c": WHITE, "b": True, "lh": 1.1}])
        yy += 0.62
    textbox(s, x + 0.4, 5.78, cw - 0.8, 0.6,
            [{"t": note, "s": 13, "c": (REDLT if main else DIM), "i": True, "lh": 1.2}])
footer(s, 13)

# ============================================================ 14. CONCLUSION
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=BG)
rect(s, 5.4, 2.05, 2.5, 0.09, fill=RED)
textbox(s, 1, 1.0, 11.33, 1.0, [{"t": "THANK YOU", "s": 15, "c": REDLT, "b": True, "sp": 4, "align": PP_ALIGN.CENTER}])
textbox(s, 1, 2.35, 11.33, 1.4, [
    {"t": "MoodFlix", "s": 60, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER, "sa": 4},
    {"t": "AI that turns a feeling into the right movie.", "s": 20, "c": GRAY, "align": PP_ALIGN.CENTER},
])
textbox(s, 1, 4.35, 11.33, 1.2, [
    {"t": "Two AI models, one 96k-row dataset, 16 moods,\nand an answer that explains itself.", "s": 16, "c": WHITE, "align": PP_ALIGN.CENTER, "lh": 1.3},
])
rect(s, 3.4, 5.95, 6.5, 0.018, fill=LINE)
textbox(s, 1, 6.15, 11.33, 0.9, [
    {"t": "Muhammad Usman · 230966      Gohar Husain · 230930      Ahmad Amir · 230948", "s": 13.5, "c": WHITE, "b": True, "align": PP_ALIGN.CENTER, "sa": 4},
    {"t": "BS IT — VI B      ·      Submitted to Ms. Warda Aslam", "s": 12, "c": GRAY, "align": PP_ALIGN.CENTER},
])

prs.save(OUT)
print("SAVED:", OUT, "slides:", len(prs.slides._sldIdLst))
